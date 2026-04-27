import pdfplumber
from pptx import Presentation
from openpyxl import load_workbook
from docx import Document
import io
from typing import Dict, Any

def parse_pdf(file_bytes: bytes) -> str:
    text = ""
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
    return text

def parse_pptx(file_bytes: bytes) -> str:
    text = ""
    prs = Presentation(io.BytesIO(file_bytes))
    for i, slide in enumerate(prs.slides):
        text += f"--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_xlsx(file_bytes: bytes) -> str:
    text = ""
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    for sheet in wb.worksheets:
        text += f"--- Sheet: {sheet.title} ---\n"
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                text += " | ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

def parse_docx(file_bytes: bytes) -> str:
    text = ""
    doc = Document(io.BytesIO(file_bytes))
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def parse_document(filename: str, content: bytes) -> str:
    ext = filename.split(".")[-1].lower()
    if ext == "pdf":
        return parse_pdf(content)
    elif ext == "pptx":
        return parse_pptx(content)
    elif ext == "xlsx":
        return parse_xlsx(content)
    elif ext == "docx":
        return parse_docx(content)
    else:
        return content.decode("utf-8", errors="ignore")
